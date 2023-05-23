import collections
import collections.abc
from pptx import Presentation


def get_text_from_slide(slide):
    """
    Get the text from a slide.
    :param slide:
    :return List[str]: A list of strings containing the text from the slide.
    """
    text_runs = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                clean_text = run.text.replace('\xa0', '')
                text_runs.append(clean_text)
    return text_runs


class Parser:
    """
    A class that parses PowerPoint presentation file and provides methods to get the text from the slides.
    """

    def __init__(self, path):
        """
            Initialize a new instance of the class with the given path.

            :param
                path (str): The path to the presentation file.
            :return: None
        """
        self.path = path
        self.presentation = Presentation(path)
        self.slides = self.presentation.slides

    def __iter__(self):
        return self.slides.__iter__()



